# reports.py — Génération de rapports PowerPoint avec visuels (charte Aid'aux)

import io
import os
from datetime import datetime

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.ticker import MaxNLocator

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

import database as db

# ── Charte graphique ──
NAVY   = RGBColor(0x00, 0x0F, 0x9F)
TEAL   = RGBColor(0x00, 0xBB, 0xB4)
ORANGE = RGBColor(0xFF, 0x47, 0x13)
RED    = RGBColor(0xC0, 0x39, 0x2B)
AMBER  = RGBColor(0xE6, 0x7E, 0x22)
YELLOW = RGBColor(0xD4, 0xAC, 0x0D)
BLUE   = RGBColor(0x29, 0x80, 0xB9)
DARK   = RGBColor(0x1A, 0x1A, 0x1A)
GREY   = RGBColor(0x6B, 0x72, 0x80)
LIGHT  = RGBColor(0xF4, 0xF6, 0xF9)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

# Palette graphiques (charte + alertes adoucies)
C_MANQ  = "#C26259"  # rouge doux (alerte)
C_PART  = "#FF4713"  # orange charte
C_COURT = "#F39C12"  # ambre
C_LONG  = "#000F9F"  # navy charte
C_OK    = "#6FB89B"  # vert doux (sans problème)
C_NA    = "#9AA3AF"  # gris (non renseigné / N/A)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

MOIS_FR = ["janvier","février","mars","avril","mai","juin",
           "juillet","août","septembre","octobre","novembre","décembre"]

BASE = os.path.dirname(os.path.abspath(__file__))
LOGO = os.path.join(BASE, "static", "logo.svg")  # SVG non supporté par pptx → on ignore
LOGO_PNG = os.path.join(BASE, "static", "icon_preview.png")


def mois_label(m):
    if not m:
        return ""
    y, mo = m.split("-")
    return f"{MOIS_FR[int(mo)-1]} {y}"


# ─────────────────────────────────────────────────────────
# Graphiques matplotlib → flux PNG
# ─────────────────────────────────────────────────────────

def _fig_to_stream(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    buf.seek(0)
    return buf


def chart_evolution(evolution):
    labels = [mois_label(e["mois"]) for e in evolution]
    fig, ax = plt.subplots(figsize=(7, 3.4))
    x = range(len(labels))
    bottom = [0] * len(labels)
    na_vals = [e.get("na", 0) for e in evolution]
    ok_vals = [max(0, e["total"] - e["manquees"] - e["partiels"] - e["courtes"] - e["longues"] - e.get("na", 0))
               for e in evolution]
    series = [
        ("Manquées", [e["manquees"] for e in evolution], C_MANQ),
        ("Badgeages partiels", [e["partiels"] for e in evolution], C_PART),
        ("Trop courtes", [e["courtes"] for e in evolution], C_COURT),
        ("Trop longues", [e["longues"] for e in evolution], C_LONG),
        ("Non renseigné", na_vals, C_NA),
        ("Sans problème", ok_vals, C_OK),
    ]
    for name, vals, color in series:
        ax.bar(x, vals, bottom=bottom, label=name, color=color, width=0.6)
        bottom = [b + v for b, v in zip(bottom, vals)]
    ax.set_xticks(list(x))
    ax.set_xticklabels(labels, rotation=0, fontsize=9)
    ax.legend(fontsize=8, loc="lower center", bbox_to_anchor=(0.5, 1.02),
              ncol=3, frameon=False, columnspacing=1.0, handletextpad=0.4)
    ax.spines[["top", "right"]].set_visible(False)
    ax.yaxis.set_major_locator(MaxNLocator(integer=True))
    ax.set_ylabel("Interventions", fontsize=9)
    return _fig_to_stream(fig)


def chart_repartition(stats, legend=False):
    probs = stats["manquees"] + stats["partiels"] + stats["courtes"] + stats["longues"]
    na = stats.get("na", 0)
    ok = max(0, stats.get("total", probs) - probs - na)
    vals = [stats["manquees"], stats["partiels"], stats["courtes"], stats["longues"], na, ok]
    labels = ["Manquées", "Badgeages\npartiels", "Trop courtes", "Trop longues", "Non renseigné", "Sans problème"]
    colors = [C_MANQ, C_PART, C_COURT, C_LONG, C_NA, C_OK]
    # Filtrer les zéros
    data = [(l, v, c) for l, v, c in zip(labels, vals, colors) if v > 0]
    fig, ax = plt.subplots(figsize=(4.4, 4.0) if legend else (5.0, 3.6))
    if data:
        ls, vs, cs = zip(*data)

        def autopct(pct):
            # Masquer le % des parts trop petites (sinon ça se chevauche)
            return f"{pct:.0f}%" if pct >= 4 else ""

        wedges, texts, autotexts = ax.pie(
            vs, colors=cs, autopct=autopct, startangle=90,
            labels=None if legend else ls,
            pctdistance=0.74 if legend else 0.78,
            labeldistance=1.13,
            textprops={"fontsize": 9},
            wedgeprops={"width": 0.45, "edgecolor": "white", "linewidth": 1.5},
        )
        for t in autotexts:
            t.set_color("white"); t.set_fontsize(9.5); t.set_fontweight("bold")

        if legend:
            # Légende sous le donut : nom + nombre (évite tout chevauchement)
            leg = [f"{l.replace(chr(10), ' ')} ({v})" for l, v in zip(ls, vs)]
            ax.legend(wedges, leg, loc="upper center", bbox_to_anchor=(0.5, -0.02),
                      ncol=2, fontsize=8.5, frameon=False, handletextpad=0.5,
                      columnspacing=1.2)
    ax.set_aspect("equal")
    return _fig_to_stream(fig)


def chart_top_intervenants(stats_list, n=10, min_interv=15):
    # On ne classe que les intervenants ayant assez d'interventions (sinon 100% sur 1 = trompeur)
    eligibles = [s for s in stats_list if s["total"] >= min_interv and s["problemes"] > 0]
    eligibles.sort(key=lambda s: s["taux"], reverse=True)
    top = eligibles[:n][::-1]  # plus haut en haut
    noms = [s["intervenant"] for s in top]
    taux = [s["taux"] for s in top]
    fig, ax = plt.subplots(figsize=(7.2, max(3, len(top) * 0.5)))
    bars = ax.barh(noms, taux, color="#000F9F", height=0.62)
    xmax = max(taux) if taux else 100
    for i, s in enumerate(top):
        ax.text(s["taux"] + xmax * 0.01, i, f"{s['taux']}%  ({s['problemes']}/{s['total']})",
                va="center", fontsize=8, color="#1A1A1A")
    ax.set_xlim(0, xmax * 1.18)
    ax.set_xlabel("Taux d'interventions à problèmes", fontsize=9)
    ax.spines[["top", "right"]].set_visible(False)
    ax.tick_params(labelsize=8)
    return _fig_to_stream(fig)


def chart_intervenant_evolution(evolution):
    labels = [mois_label(e["mois"]) for e in evolution]
    fig, ax = plt.subplots(figsize=(6.0, 2.7))
    x = range(len(labels))
    bottom = [0] * len(labels)
    ok_vals = [max(0, e["total"] - e["manquees"] - e["partiels"] - e["courtes"] - e["longues"] - e.get("na", 0))
               for e in evolution]
    for name, vals, color in [("Manquées",[e["manquees"] for e in evolution],C_MANQ),
                              ("Badgeages partiels",[e["partiels"] for e in evolution],C_PART),
                              ("Trop courtes",[e["courtes"] for e in evolution],C_COURT),
                              ("Trop longues",[e["longues"] for e in evolution],C_LONG),
                              ("Non renseigné",[e.get("na",0) for e in evolution],C_NA),
                              ("Sans problème",ok_vals,C_OK)]:
        ax.bar(x, vals, bottom=bottom, label=name, color=color, width=0.6)
        bottom = [b + v for b, v in zip(bottom, vals)]
    ax.set_xticks(list(x)); ax.set_xticklabels(labels, fontsize=9)
    ax.legend(fontsize=7, loc="lower center", bbox_to_anchor=(0.5, 1.02),
              ncol=3, frameon=False, columnspacing=0.9, handletextpad=0.4)
    ax.spines[["top","right"]].set_visible(False)
    ax.yaxis.set_major_locator(MaxNLocator(integer=True))
    return _fig_to_stream(fig)


# ─────────────────────────────────────────────────────────
# Helpers PPTX
# ─────────────────────────────────────────────────────────

def _add_textbox(slide, left, top, width, height, text, size=18, bold=False,
                 color=DARK, align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    return box


def _add_band(slide, color, top=Inches(0), height=Inches(1.3)):
    band = slide.shapes.add_shape(1, Inches(0), top, SLIDE_W, height)  # rectangle
    band.fill.solid()
    band.fill.fore_color.rgb = color
    band.line.fill.background()
    band.shadow.inherit = False
    return band


def _kpi_box(slide, left, top, value, label, color):
    w, h = Inches(2.7), Inches(1.5)
    card = slide.shapes.add_shape(1, left, top, w, h)
    card.fill.solid(); card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(0xE5, 0xE7, 0xEB); card.line.width = Pt(0.75)
    card.shadow.inherit = False

    tf = card.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p1 = tf.paragraphs[0]; p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run(); r1.text = str(value)
    r1.font.size = Pt(34); r1.font.bold = True; r1.font.color.rgb = color
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = label
    r2.font.size = Pt(11); r2.font.color.rgb = GREY


def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ─────────────────────────────────────────────────────────
# Slides
# ─────────────────────────────────────────────────────────

def _oval(s, left, top, w, h, color):
    o = s.shapes.add_shape(9, left, top, w, h)  # OVAL
    o.fill.solid(); o.fill.fore_color.rgb = color
    o.line.fill.background(); o.shadow.inherit = False
    return o


def _rect(s, left, top, w, h, color, rounded=False):
    o = s.shapes.add_shape(5 if rounded else 1, left, top, w, h)
    o.fill.solid(); o.fill.fore_color.rgb = color
    o.line.fill.background(); o.shadow.inherit = False
    return o


def _tricolor_strip(s, top=Inches(0), height=Inches(0.13)):
    """Fine bande tricolore (navy / teal / orange) en haut de la diapo."""
    w = SLIDE_W
    _rect(s, Inches(0), top, Emu(int(w * 0.55)), height, NAVY)
    _rect(s, Emu(int(w * 0.55)), top, Emu(int(w * 0.30)), height, TEAL)
    _rect(s, Emu(int(w * 0.85)), top, Emu(int(w * 0.15)), height, ORANGE)


def slide_titre(prs, periode):
    s = _blank_slide(prs)
    # Fond blanc
    bg = _rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, WHITE)

    # Formes décoratives aux couleurs de la charte (en grande partie hors-cadre)
    _oval(s, Inches(10.6), Inches(-2.0), Inches(5), Inches(5), TEAL)      # haut-droite
    _oval(s, Inches(-1.8), Inches(5.2), Inches(4), Inches(4), ORANGE)     # bas-gauche
    _oval(s, Inches(11.6), Inches(5.6), Inches(2.4), Inches(2.4), NAVY)   # bas-droite
    # Bande tricolore en haut
    _tricolor_strip(s)

    # Logo (fond blanc → pas besoin de carton)
    if os.path.exists(LOGO_PNG):
        try:
            s.shapes.add_picture(LOGO_PNG, Inches(5.46), Inches(1.5), height=Inches(2.0))
        except Exception:
            pass

    _add_textbox(s, Inches(1), Inches(3.9), Inches(11.3), Inches(1),
                 "Rapport de suivi des interventions", size=36, bold=True,
                 color=NAVY, align=PP_ALIGN.CENTER)
    _add_textbox(s, Inches(1), Inches(4.9), Inches(11.3), Inches(0.6),
                 periode, size=18, color=TEAL, align=PP_ALIGN.CENTER)

    # Petit accent tricolore sous le sous-titre
    cx = SLIDE_W / 2
    bw = Inches(0.7)
    _rect(s, Emu(int(cx - bw * 1.6)), Inches(5.55), bw, Inches(0.09), NAVY, rounded=True)
    _rect(s, Emu(int(cx - bw * 0.5)), Inches(5.55), bw, Inches(0.09), TEAL, rounded=True)
    _rect(s, Emu(int(cx + bw * 0.6)), Inches(5.55), bw, Inches(0.09), ORANGE, rounded=True)

    _add_textbox(s, Inches(1), Inches(6.7), Inches(11.3), Inches(0.5),
                 f"Généré le {datetime.now().strftime('%d/%m/%Y')}  ·  Aid'aux — Services d'aide à la personne",
                 size=11, color=GREY, align=PP_ALIGN.CENTER)
    return s


def _entete(s, titre, sous_titre=""):
    """En-tête sur fond blanc : bande tricolore + titre navy + accent latéral."""
    _tricolor_strip(s)
    # Petit bloc d'accent à gauche du titre
    _rect(s, Inches(0.6), Inches(0.42), Inches(0.14), Inches(0.62), TEAL, rounded=True)
    _add_textbox(s, Inches(0.95), Inches(0.32), Inches(11.5), Inches(0.6),
                 titre, size=26, bold=True, color=NAVY)
    if sous_titre:
        _add_textbox(s, Inches(0.97), Inches(0.88), Inches(11.5), Inches(0.4),
                     sous_titre, size=13, color=GREY)
    # Ligne de séparation fine sous l'en-tête
    _rect(s, Inches(0.6), Inches(1.32), Inches(12.13), Pt(1.2), RGBColor(0xE0,0xE3,0xEA))


def slide_synthese(prs, info, stats, evolution, periode):
    s = _blank_slide(prs)
    _entete(s, "Synthèse globale", periode)

    total = stats["total"] or 1
    probs = stats["manquees"] + stats["partiels"] + stats["courtes"] + stats["longues"]

    # KPIs
    y = Inches(1.4)
    _kpi_box(s, Inches(0.6), y, info["total"], "Interventions totales", NAVY)
    _kpi_box(s, Inches(3.5), y, stats["manquees"],
             f"Manquées ({round(stats['manquees']/total*100,1)}%)", RED)
    _kpi_box(s, Inches(6.4), y, stats["partiels"],
             f"Badgeages partiels ({round(stats['partiels']/total*100,1)}%)", AMBER)
    _kpi_box(s, Inches(9.3), y, stats["courtes"] + stats["longues"],
             f"Durée incorrecte ({round((stats['courtes']+stats['longues'])/total*100,1)}%)", BLUE)

    # Graphiques
    if evolution:
        img = chart_evolution(evolution)
        s.shapes.add_picture(img, Inches(0.5), Inches(3.2), width=Inches(7.6))
    img2 = chart_repartition(stats)
    s.shapes.add_picture(img2, Inches(8.4), Inches(3.2), width=Inches(4.4))

    _add_textbox(s, Inches(0.6), Inches(3.0), Inches(7), Inches(0.4),
                 "Évolution mensuelle", size=13, bold=True, color=NAVY)
    _add_textbox(s, Inches(8.5), Inches(3.0), Inches(4), Inches(0.4),
                 "Répartition par type", size=13, bold=True, color=NAVY)
    return s


def slide_classements(prs, stats_list, top_clients, periode):
    s = _blank_slide(prs)
    _entete(s, "Classements & alertes", periode)

    # Graphique top intervenants
    _add_textbox(s, Inches(0.6), Inches(1.3), Inches(7), Inches(0.4),
                 "Intervenants les plus en difficulté", size=14, bold=True, color=NAVY)
    _add_textbox(s, Inches(0.62), Inches(1.68), Inches(7), Inches(0.3),
                 "Parmi les intervenants ayant au moins 15 interventions sur la période",
                 size=10, color=GREY)
    img = chart_top_intervenants(stats_list, n=10)
    s.shapes.add_picture(img, Inches(0.5), Inches(2.1), width=Inches(7.4))

    # Tableau top clients
    _add_textbox(s, Inches(8.4), Inches(1.3), Inches(4.5), Inches(0.4),
                 "Bénéficiaires les plus impactés", size=14, bold=True, color=NAVY)
    rows = len(top_clients) + 1
    tbl = s.shapes.add_table(rows, 2, Inches(8.4), Inches(1.8),
                             Inches(4.4), Inches(0.4 * rows)).table
    tbl.cell(0, 0).text = "Bénéficiaire"
    tbl.cell(0, 1).text = "Problèmes"
    for j in range(2):
        c = tbl.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = NAVY
        c.text_frame.paragraphs[0].runs[0].font.color.rgb = WHITE
        c.text_frame.paragraphs[0].runs[0].font.bold = True
        c.text_frame.paragraphs[0].runs[0].font.size = Pt(11)
    for i, cl in enumerate(top_clients, start=1):
        tbl.cell(i, 0).text = cl["client"]
        tbl.cell(i, 1).text = str(cl["problemes"])
        for j in range(2):
            tbl.cell(i, j).text_frame.paragraphs[0].runs[0].font.size = Pt(10)
    return s


def slide_intervenant(prs, nom, detail, moyenne_taux):
    s = _blank_slide(prs)
    g = detail["global"]
    _entete(s, nom, "Rapport individuel")

    y = Inches(1.4)
    _kpi_box(s, Inches(0.6), y, g["total"], "Interventions", NAVY)
    _kpi_box(s, Inches(3.5), y, g["problemes"], f"À problèmes ({g['taux']}%)",
             RED if g["taux"] > moyenne_taux else TEAL)
    _kpi_box(s, Inches(6.4), y, g["manquees"], "Manquées", RED)
    _kpi_box(s, Inches(9.3), y, g["nb_clients"], "Bénéficiaires suivis", NAVY)

    # Comparaison à la moyenne
    diff = round(g["taux"] - moyenne_taux, 1)
    if diff > 0:
        txt = f"⚠ Taux supérieur de {diff} pts à la moyenne ({moyenne_taux}%)"
        col = RED
    else:
        txt = f"✓ Taux inférieur de {abs(diff)} pts à la moyenne ({moyenne_taux}%)"
        col = RGBColor(0x1E, 0x88, 0x49)
    _add_textbox(s, Inches(0.6), Inches(3.05), Inches(12), Inches(0.4),
                 txt, size=14, bold=True, color=col)

    # Détail par type
    _add_textbox(s, Inches(0.6), Inches(3.6), Inches(12), Inches(0.3),
                 f"Manquées : {g['manquees']}   ·   Badgeages partiels : {g['partiels']}   ·   "
                 f"Trop courtes : {g['courtes']}   ·   Trop longues : {g['longues']}"
                 + (f"   ·   Non renseigné : {g.get('na', 0)}" if g.get('na') else ""),
                 size=12, color=DARK)

    evo = detail["evolution"]
    if len(evo) > 1:
        # Évolution (gauche) + répartition (droite)
        _add_textbox(s, Inches(0.6), Inches(4.05), Inches(6), Inches(0.3),
                     "Évolution mensuelle", size=13, bold=True, color=NAVY)
        img = chart_intervenant_evolution(evo)
        s.shapes.add_picture(img, Inches(0.4), Inches(4.45), height=Inches(2.75))

        _add_textbox(s, Inches(8.9), Inches(4.05), Inches(4), Inches(0.3),
                     "Répartition", size=13, bold=True, color=NAVY)
        img2 = chart_repartition(g, legend=True)
        s.shapes.add_picture(img2, Inches(8.7), Inches(4.4), height=Inches(2.7))
    else:
        # Un seul mois : camembert centré
        _add_textbox(s, Inches(0.6), Inches(4.05), Inches(6), Inches(0.3),
                     "Répartition", size=13, bold=True, color=NAVY)
        img2 = chart_repartition(g, legend=True)
        s.shapes.add_picture(img2, Inches(4.6), Inches(4.3), height=Inches(2.8))
    return s


def slide_mensuel(prs, mois_data):
    s = _blank_slide(prs)
    _entete(s, mois_label(mois_data["mois"]), "Bilan mensuel")

    total = mois_data["total"] or 1
    y = Inches(1.4)
    _kpi_box(s, Inches(0.6), y, mois_data["total"], "Interventions", NAVY)
    _kpi_box(s, Inches(3.5), y, mois_data["problemes"],
             f"À problèmes ({mois_data['taux']}%)", AMBER)
    _kpi_box(s, Inches(6.4), y, mois_data["manquees"], "Manquées", RED)
    _kpi_box(s, Inches(9.3), y, mois_data["courtes"] + mois_data["longues"],
             "Durée incorrecte", BLUE)

    # Variation vs mois précédent
    if mois_data.get("variation") is not None:
        v = mois_data["variation"]
        if v > 0:
            txt, col = f"▲ Taux de problèmes en hausse de {v} pts vs mois précédent", RED
        elif v < 0:
            txt, col = f"▼ Taux de problèmes en baisse de {abs(v)} pts vs mois précédent", RGBColor(0x1E,0x88,0x49)
        else:
            txt, col = "= Taux de problèmes stable vs mois précédent", GREY
    else:
        txt, col = "Premier mois de la période — pas de comparaison disponible", GREY
    _add_textbox(s, Inches(0.6), Inches(3.2), Inches(12), Inches(0.5),
                 txt, size=15, bold=True, color=col)

    # Détail par type (colonne gauche)
    _add_textbox(s, Inches(0.6), Inches(4.1), Inches(5), Inches(0.4),
                 "Détail par type de problème", size=14, bold=True, color=NAVY)
    details = [
        ("Manquées", mois_data["manquees"], C_MANQ),
        ("Badgeages partiels", mois_data["partiels"], C_PART),
        ("Trop courtes", mois_data["courtes"], C_COURT),
        ("Trop longues", mois_data["longues"], C_LONG),
        ("Non renseigné", mois_data.get("na", 0), C_NA),
    ]
    yy = 4.7
    for label, val, hexcol in details:
        rgb = RGBColor(int(hexcol[1:3],16), int(hexcol[3:5],16), int(hexcol[5:7],16))
        dot = s.shapes.add_shape(9, Inches(0.7), Inches(yy+0.05), Inches(0.18), Inches(0.18))
        dot.fill.solid(); dot.fill.fore_color.rgb = rgb; dot.line.fill.background(); dot.shadow.inherit=False
        _add_textbox(s, Inches(1.0), Inches(yy-0.08), Inches(3), Inches(0.4),
                     f"{label} : {val}", size=13, color=DARK)
        yy += 0.5

    # Camembert (droite)
    img = chart_repartition({
        "total": mois_data["total"], "na": mois_data.get("na", 0),
        "manquees": mois_data["manquees"], "partiels": mois_data["partiels"],
        "courtes": mois_data["courtes"], "longues": mois_data["longues"],
    })
    s.shapes.add_picture(img, Inches(7.4), Inches(3.9), width=Inches(5.0))
    return s


# ─────────────────────────────────────────────────────────
# Assemblage
# ─────────────────────────────────────────────────────────

def generer_rapport(sections, mois=None, intervenants=None):
    """
    sections : liste parmi ['synthese','classements','intervenants','mensuel']
    mois : filtre optionnel (YYYY-MM) pour la synthèse/classements
    intervenants : liste optionnelle de noms. Si fournie, les fiches ne concernent
                   que ces intervenants (même sans problème). Sinon, tous ceux
                   ayant au moins un problème.
    Retourne un BytesIO du .pptx.
    """
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    info = db.get_global_info()
    periode = ""
    if info["period_min"] and info["period_max"]:
        def fr(iso):
            d = iso.split(" ")[0].split("-")
            return f"{d[2]}/{d[1]}/{d[0]}"
        periode = f"Période : {fr(info['period_min'])} – {fr(info['period_max'])}"
    if mois:
        periode = f"Mois : {mois_label(mois)}"

    slide_titre(prs, periode)

    stats = db.get_stats(mois=mois)
    evolution = db.get_monthly_evolution()

    if "synthese" in sections:
        slide_synthese(prs, info, stats, evolution, periode)

    if "classements" in sections:
        stats_list = db.stats_par_intervenant(mois=mois)
        top_clients = db.top_clients_problemes(mois=mois, limit=10)
        slide_classements(prs, stats_list, top_clients, periode)

    if "intervenants" in sections:
        stats_list = db.stats_par_intervenant(mois=mois)
        # Moyenne globale des taux (sur tous, pour la comparaison)
        moyenne = round(sum(s["taux"] for s in stats_list) / len(stats_list), 1) if stats_list else 0

        if intervenants:
            # Sélection précise : on garde l'ordre choisi, on inclut même sans problème
            voulus = set(intervenants)
            choisis = [st for st in stats_list if st["intervenant"] in voulus]
        else:
            # Tous ceux ayant au moins un problème (triés par gravité)
            choisis = [st for st in stats_list if st["problemes"] > 0]

        for st in choisis:
            detail = db.detail_intervenant(st["intervenant"])
            slide_intervenant(prs, st["intervenant"], detail, moyenne)

    if "mensuel" in sections:
        for md in db.stats_mensuelles_detaillees():
            slide_mensuel(prs, md)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
